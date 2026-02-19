"""PowerPoint to Markdown converter.

Converts PPTX presentations to Markdown with:
- Slide title extraction
- Table conversion
- Hash-based image naming for diff stability
- YAML frontmatter and TOC generation
- Optional speaker notes
"""

from __future__ import annotations

from pathlib import Path
from typing import TYPE_CHECKING, Any

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.shapes.picture import Picture
except ImportError as e:
    raise ImportError(
        "python-pptx is required for convert. Install with: pip install python-pptx"
    ) from e

from otutil.tools._convert.utils import (
    IncrementalWriter,
    compute_file_checksum,
    get_mtime_iso,
    normalise_whitespace,
    save_image,
    write_toc_file,
)

if TYPE_CHECKING:
    from pptx.presentation import Presentation as PresentationType
    from pptx.shapes.base import BaseShape
    from pptx.slide import Slide


def convert_powerpoint(
    input_path: Path,
    output_dir: Path,
    source_rel: str,
    *,
    include_notes: bool = False,
) -> dict[str, Any]:
    """Convert PowerPoint presentation to Markdown.

    Args:
        input_path: Path to PPTX file
        output_dir: Directory for output files
        source_rel: Relative path to source for frontmatter
        include_notes: Include speaker notes after slide content

    Returns:
        Dict with 'output', 'slides', 'images' keys
    """
    output_dir.mkdir(parents=True, exist_ok=True)

    prs: PresentationType = Presentation(str(input_path))
    try:
        # Get metadata for frontmatter
        checksum = compute_file_checksum(input_path)
        mtime = get_mtime_iso(input_path)
        total_slides = len(prs.slides)

        # Set up images directory
        images_dir = output_dir / f"{input_path.stem}_images"
        writer = IncrementalWriter()
        images_extracted = 0

        # Process slides
        for slide_idx, slide in enumerate(prs.slides, 1):
            imgs = _process_slide(
                slide, slide_idx, writer, images_dir, include_notes
            )
            images_extracted += imgs
    finally:
        # Ensure presentation resources are released
        # python-pptx Presentation doesn't have explicit close, but we can
        # help garbage collection by clearing references
        del prs

    # Write main output (pure content, no frontmatter - line numbers start at 1)
    content = normalise_whitespace(writer.get_content())
    output_path = output_dir / f"{input_path.stem}.md"
    output_path.write_text(content, encoding="utf-8")

    # Write separate TOC file (includes frontmatter)
    headings = writer.get_headings()
    toc_path = write_toc_file(
        headings=headings,
        output_dir=output_dir,
        stem=input_path.stem,
        source=source_rel,
        converted=mtime,
        pages=total_slides,
        checksum=checksum,
    )

    return {
        "output": str(output_path),
        "toc": str(toc_path),
        "slides": total_slides,
        "images": images_extracted,
    }


def _process_slide(
    slide: Slide,
    slide_number: int,
    writer: IncrementalWriter,
    images_dir: Path,
    include_notes: bool,
) -> int:
    """Process a single slide.

    Returns:
        Number of images extracted
    """
    images_extracted = 0

    # Get slide title
    title = f"Slide {slide_number}"
    try:
        if (
            hasattr(slide, "shapes")
            and hasattr(slide.shapes, "title")
            and slide.shapes.title
            and slide.shapes.title.text.strip()
        ):
            title = slide.shapes.title.text.strip()
    except AttributeError:
        pass

    writer.write_heading(2, title)

    # Process shapes
    text_content: list[str] = []
    tables_content: list[str] = []

    for shape in slide.shapes:
        shape_type = getattr(shape, "shape_type", None)

        # Tables (check FIRST - tables also have .text attribute)
        if shape_type == MSO_SHAPE_TYPE.TABLE:
            try:
                if hasattr(shape, "table"):
                    table_md = _process_table(shape.table)
                    if table_md:
                        tables_content.append(table_md)
            except (AttributeError, ValueError):
                pass

        # Images
        elif shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                if isinstance(shape, Picture):
                    img_ref = _process_image(shape, images_dir)
                    if img_ref:
                        writer.write(img_ref + "\n\n")
                        images_extracted += 1
            except Exception:
                pass

        # Text shapes (check LAST - after tables and images)
        elif hasattr(shape, "text") and shape.text and shape.text.strip():
            # Skip title shape (already processed)
            if shape == getattr(slide.shapes, "title", None):
                continue
            text_content.append(_process_text_shape(shape))

    # Write text content
    for text in text_content:
        if text:
            writer.write(text + "\n\n")

    # Write tables
    for table in tables_content:
        if table:
            writer.write(table + "\n\n")

    # Add speaker notes if requested
    if include_notes and hasattr(slide, "notes_slide"):
        try:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame and notes_frame.text.strip():
                writer.write("**Speaker Notes:**\n\n")
                writer.write(f"> {notes_frame.text.strip()}\n\n")
        except Exception:
            pass

    # Add slide separator
    writer.write("---\n\n")

    return images_extracted


def _process_text_shape(shape: BaseShape) -> str:
    """Process text from a shape."""
    if not hasattr(shape, "text") or not shape.text.strip():
        return ""

    text: str = str(shape.text).strip()
    lines = text.split("\n")

    if len(lines) > 1:
        # Format as bullet list
        processed: list[str] = []
        for line in lines:
            line = line.strip()
            if line:
                # Remove existing bullet markers
                for marker in ("•", "-", "*", "○", "▪", "▫"):
                    if line.startswith(marker):
                        line = line[1:].strip()
                        break
                processed.append(f"- {line}")
        return "\n".join(processed)

    return text


def _process_table(table: Any) -> str:
    """Convert table to Markdown."""
    if not hasattr(table, "rows") or not table.rows:
        return ""

    lines: list[str] = []

    # Header row
    header_cells: list[str] = []
    for cell in table.rows[0].cells:
        header_cells.append(cell.text.strip() if hasattr(cell, "text") else "")

    if not header_cells:
        return ""

    lines.append("| " + " | ".join(header_cells) + " |")
    lines.append("| " + " | ".join("---" for _ in header_cells) + " |")

    # Data rows
    for i in range(1, len(table.rows)):
        row = table.rows[i]
        cells: list[str] = []
        for cell in row.cells:
            cells.append(cell.text.strip() if hasattr(cell, "text") else "")

        while len(cells) < len(header_cells):
            cells.append("")

        lines.append("| " + " | ".join(cells[: len(header_cells)]) + " |")

    return "\n".join(lines)


def _process_image(shape: Picture, images_dir: Path) -> str:
    """Extract and save image with hash-based naming."""
    try:
        image_data = shape.image.blob
        content_type = shape.image.content_type

        img_path = save_image(image_data, images_dir, content_type)
        rel_path = f"{images_dir.name}/{img_path.name}"
        return f"![{img_path.name}]({rel_path})"

    except Exception:
        return ""
