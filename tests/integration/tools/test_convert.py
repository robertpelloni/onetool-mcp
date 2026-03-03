"""Integration tests for document conversion tools.

Creates real documents using each library and verifies the conversion
produces meaningful Markdown output.
"""

from __future__ import annotations

import os
from pathlib import Path
from typing import Any

import pytest


@pytest.fixture
def output_dir(tmp_path: Path) -> Path:
    """Temporary directory for conversion output files."""
    d = tmp_path / "output"
    d.mkdir()
    return d


@pytest.mark.integration
@pytest.mark.tools
class TestConvertPdf:
    """Integration tests for convert.pdf()."""

    def test_pdf_to_markdown(self, tmp_path: Path, output_dir: Path) -> None:
        """Create a minimal PDF with fitz and verify text is extracted."""
        try:
            import fitz
        except ImportError:
            pytest.fail("pymupdf not installed")

        doc: Any = fitz.open()
        page = doc.new_page()
        page.insert_text((50, 100), "Hello PDF World")
        pdf_path = tmp_path / "test.pdf"
        doc.save(str(pdf_path))
        doc.close()

        from otutil.tools.convert import pdf

        old_cwd = os.getcwd()
        try:
            os.chdir(tmp_path)
            result = pdf(pattern=str(pdf_path), output_dir=str(output_dir))
        finally:
            os.chdir(old_cwd)

        assert "test.pdf" in result
        output_files = list(output_dir.glob("*.md"))
        assert len(output_files) >= 1
        assert "Hello PDF World" in "\n".join(f.read_text() for f in output_files)


@pytest.mark.integration
@pytest.mark.tools
class TestConvertWord:
    """Integration tests for convert.word()."""

    def test_word_to_markdown(self, tmp_path: Path, output_dir: Path) -> None:
        """Create a DOCX with python-docx and verify content is extracted."""
        try:
            import docx
        except ImportError:
            pytest.fail("python-docx not installed")

        doc = docx.Document()
        doc.add_heading("Test Document", level=1)
        doc.add_paragraph("This is the first paragraph.")
        docx_path = tmp_path / "test.docx"
        doc.save(str(docx_path))

        from otutil.tools.convert import word

        old_cwd = os.getcwd()
        try:
            os.chdir(tmp_path)
            result = word(pattern=str(docx_path), output_dir=str(output_dir))
        finally:
            os.chdir(old_cwd)

        assert "test.docx" in result
        output_files = list(output_dir.glob("*.md"))
        assert len(output_files) >= 1
        all_content = "\n".join(f.read_text() for f in output_files)
        assert "Test Document" in all_content or "first paragraph" in all_content


@pytest.mark.integration
@pytest.mark.tools
class TestConvertPptx:
    """Integration tests for convert.powerpoint()."""

    def test_pptx_to_markdown(self, tmp_path: Path, output_dir: Path) -> None:
        """Create a PPTX with python-pptx and verify slides are extracted."""
        try:
            import pptx
        except ImportError:
            pytest.fail("python-pptx not installed")

        prs = pptx.Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        if title:
            title.text = "Slide One Title"
        pptx_path = tmp_path / "test.pptx"
        prs.save(str(pptx_path))

        from otutil.tools.convert import powerpoint

        old_cwd = os.getcwd()
        try:
            os.chdir(tmp_path)
            result = powerpoint(pattern=str(pptx_path), output_dir=str(output_dir))
        finally:
            os.chdir(old_cwd)

        assert "test.pptx" in result
        output_files = list(output_dir.glob("*.md"))
        assert len(output_files) >= 1


@pytest.mark.integration
@pytest.mark.tools
class TestConvertExcel:
    """Integration tests for convert.excel()."""

    def test_excel_to_markdown(self, tmp_path: Path, output_dir: Path) -> None:
        """Create an XLSX with openpyxl and verify sheet data is extracted."""
        try:
            import openpyxl
        except ImportError:
            pytest.fail("openpyxl not installed")

        wb = openpyxl.Workbook()
        ws = wb.active
        assert ws is not None
        ws["A1"] = "Name"
        ws["B1"] = "Score"
        ws["A2"] = "Alice"
        ws["B2"] = 95
        xlsx_path = tmp_path / "test.xlsx"
        wb.save(str(xlsx_path))
        wb.close()

        from otutil.tools.convert import excel

        old_cwd = os.getcwd()
        try:
            os.chdir(tmp_path)
            result = excel(pattern=str(xlsx_path), output_dir=str(output_dir))
        finally:
            os.chdir(old_cwd)

        assert "test.xlsx" in result
        output_files = list(output_dir.glob("*.md"))
        assert len(output_files) >= 1
        all_content = "\n".join(f.read_text() for f in output_files)
        assert "Alice" in all_content or "Name" in all_content
